from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

def add_box(slide, left, top, width, height, text, bg_color, font_color=RGBColor(0,0,0), font_size=14, bold=False, sub_text=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = RGBColor(80, 80, 80)
    shape.line.width = Pt(1.2)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # vertical center
    tf.vertical_anchor = 3  # middle

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color

    if sub_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub_text
        run2.font.size = Pt(10)
        run2.font.color.rgb = font_color

    return shape

def add_arrow(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = RGBColor(60, 60, 60)
    connector.line.width = Pt(1.5)

def add_text(slide, left, top, width, height, text, font_size=10, color=RGBColor(80,80,80), bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold

# ── 타이틀 ──────────────────────────────────────
add_text(slide, 0.3, 0.1, 12.5, 0.5, "Doosan Robot2 ROS2 Software Block Diagram",
         font_size=20, color=RGBColor(20,60,120), bold=True)

# ── 색상 정의 ──────────────────────────────────
C_USER    = RGBColor(198, 224, 180)   # 연두
C_ROS     = RGBColor(189, 215, 238)   # 연파랑
C_CTRL    = RGBColor(255, 230, 153)   # 연노랑
C_LIB     = RGBColor(255, 192, 0)     # 주황노랑
C_HW      = RGBColor(217, 150, 148)   # 연빨강
C_MSG     = RGBColor(228, 208, 255)   # 연보라
C_GRAY    = RGBColor(230, 230, 230)

WHITE     = RGBColor(255,255,255)
DARK      = RGBColor(30,30,30)

# ── 블록 레이아웃 (left, top, width, height) ─────
# 세로 중심 구조: 위→아래

# [1] User Application
add_box(slide, 0.4, 0.75, 3.5, 0.85,
        "User Application / ROS2 Node", C_USER, DARK, 13, True,
        "Python / C++ 사용자 코드")

# [2] dsr_msgs2
add_box(slide, 4.2, 0.75, 2.3, 0.85,
        "dsr_msgs2", C_MSG, DARK, 13, True,
        "Custom Service/Topic\n메시지 정의")

# [3] dsr_bringup2
add_box(slide, 0.4, 2.1, 2.2, 0.85,
        "dsr_bringup2", C_ROS, DARK, 13, True,
        "Launch files\n파라미터 설정")

# [4] dsr_controller2
add_box(slide, 2.9, 2.1, 2.8, 0.85,
        "dsr_controller2", C_CTRL, DARK, 13, True,
        "ros2_control\nHardware Interface")

# [5] common2 / DRFL
add_box(slide, 0.4, 3.5, 5.3, 1.0,
        "common2  (DRFL - Doosan Robot Framework Library)", C_LIB, DARK, 13, True,
        "libDRFL.a  |  TCP Socket 통신  |  로봇 제어 API")

# [6] Real Robot / Emulator
add_box(slide, 0.4, 5.1, 2.3, 0.9,
        "Doosan Robot\n(Real)", C_HW, DARK, 12, True,
        "port 12345 / 19999")

add_box(slide, 3.4, 5.1, 2.3, 0.9,
        "DSR Emulator\n(Docker)", C_HW, DARK, 12, True,
        "port 12345 / 19999")

# ── ROS2 Interface 박스 (우측 설명) ──────────────
add_box(slide, 6.8, 0.75, 5.9, 2.8,
        "ROS2 Interface", RGBColor(240,248,255), DARK, 13, True, None)

add_text(slide, 7.0, 1.1, 5.5, 2.2,
"""[ Services ]
/dsr01/motion/move_joint        관절 이동
/dsr01/motion/move_line         직선 이동
/dsr01/motion/move_circle       원호 이동
/dsr01/system/set_robot_mode    모드 설정
/dsr01/gripper/...              그리퍼 제어

[ Topics ]
/dsr01/joint_states             관절 각도/속도/토크 (rad)
/dsr01/state                    로봇 상태""",
         font_size=10.5, color=RGBColor(20,20,80))

# ── 빌드 구조 박스 ─────────────────────────────
add_box(slide, 6.8, 3.7, 5.9, 2.3,
        "Build Structure", RGBColor(245,245,245), DARK, 13, True, None)

add_text(slide, 7.0, 4.0, 5.5, 1.8,
"""~/hye/doosan_ros2_ws/src/doosan-robot2/
  ├── common2/          DRFL 라이브러리 (libDRFL.a)
  ├── dsr_bringup2/     Launch 파일
  ├── dsr_controller2/  ros2_control HW Interface
  ├── dsr_msgs2/        커스텀 메시지/서비스
  ├── dsr_tests/        테스트 코드
  └── gz_ros2_control/  Gazebo 연동""",
         font_size=10.5, color=RGBColor(20,20,80))

# ── 화살표 ─────────────────────────────────────
# User → dsr_bringup2
add_arrow(slide, 1.5, 1.6, 1.5, 2.1)
# User → dsr_controller2
add_arrow(slide, 3.0, 1.6, 3.7, 2.1)
# dsr_bringup2 → DRFL
add_arrow(slide, 1.5, 2.95, 1.5, 3.5)
# dsr_controller2 → DRFL
add_arrow(slide, 3.7, 2.95, 3.0, 3.5)
# DRFL → Real Robot
add_arrow(slide, 1.5, 4.5, 1.5, 5.1)
# DRFL → Emulator
add_arrow(slide, 4.0, 4.5, 4.5, 5.1)
# dsr_msgs2 → User
add_arrow(slide, 4.2, 1.18, 3.9, 1.18)

# ── 레이블 ────────────────────────────────────
add_text(slide, 1.55, 1.75, 1.2, 0.3, "launch", 9, RGBColor(100,100,100))
add_text(slide, 3.1, 1.75, 1.5, 0.3, "HW Interface", 9, RGBColor(100,100,100))
add_text(slide, 1.55, 3.1, 1.5, 0.3, "TCP 12345", 9, RGBColor(150,60,60))
add_text(slide, 3.55, 3.1, 1.5, 0.3, "TCP 12345", 9, RGBColor(150,60,60))

# ── 저장 ─────────────────────────────────────
out = "/home/hanku/hye/doosan_robot2_block.pptx"
prs.save(out)
print(f"Saved: {out}")
